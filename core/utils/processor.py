from django.conf import settings
import docx2txt
from PyPDF2 import PdfReader
import textract
import os
import re
from pptx import Presentation

class DocumentProcessor:
    """
    Extracts text from a document file.
    Supports .docx, .pptx, .pdf, .txt files.
    """

    def __init__(self, file_path):
        """
        Initializes the DocumentProcessor.

        Args:
            file_path (str): The path of the document file.
        """
        self.file_path = self.get_file_path(file_path)
        self.file_extension = self.get_file_path(file_path).split(".")[-1]
        self.text = ""

    def extract_text(self):
        """
        Extracts the text from the document file based on its extension.
        """
        if self.file_extension == "docx":
            self.text = docx2txt.process(self.file_path)
        elif self.file_extension == "pptx":
            ppt = Presentation(self.file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        self.text += shape.text
        elif self.file_extension == "pdf":
            pdf_reader = PdfReader(self.file_path)
            for page in pdf_reader.pages:
                self.text += page.extract_text()
        elif self.file_extension == "txt":
            with open(self.file_path, "r") as f:
                self.text = f.read()
        else:
            raise ValueError("Unsupported file extension")

    def get_text(self):
        """
        Retrieves the extracted text.

        Returns:
            str: The extracted text.
        """
        if not self.text:
            self.extract_text()
        return self.text

    def process_text(self):
        """
        Processes the extracted text to remove unwanted characters and make meaningful paragraphs.

        Returns:
            list: A list of dictionaries, where each dictionary contains the processed text as "content".
        """
        # Replace any unwanted characters with spaces
        self.text = re.sub('[^a-zA-Z0-9]', ' ', self.text)

        # Split the text into paragraphs based on common delimiters
        paragraphs = re.split('\n|\r\n*\r\n|\n\t*\n|\r\n\t*\r\n', self.text)

        # Join the paragraphs back together with double newlines
        processed_text = '\n\n'.join(paragraphs)

        # Create a list of dictionaries, where each dictionary contains a paragraph of processed text
        docs = [{"content": text} for text in paragraphs]

        return docs

    def get_file_path(self, file):
        """
        Returns the absolute file path by appending the file name to the COURSE_MATERIALS_DIR from settings.py.

        Args:
            file (str): The file path.

        Returns:
            str: The absolute file path.
        """
        return os.path.join(settings.COURSE_MATERIALS_DIR, file.split("/")[-1])
